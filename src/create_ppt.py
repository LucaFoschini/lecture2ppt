import re
from datetime import datetime, timedelta
import os
import argparse
from pptx import Presentation
from pptx.util import Inches
import fitz
from PIL import Image
import io
import json

def parse_timestamp(timestamp_str):
    """Convert timestamp string to timedelta"""
    try:
        # Handle different timestamp formats
        if '.' in timestamp_str:
            # Format: HH:MM:SS.mmm
            hours, minutes, rest = timestamp_str.split(':')
            seconds, milliseconds = rest.split('.')
            return timedelta(hours=float(hours), 
                           minutes=float(minutes), 
                           seconds=float(seconds), 
                           milliseconds=float(milliseconds))
        else:
            # Format: HH:MM:SS
            hours, minutes, seconds = map(float, timestamp_str.split(':'))
            return timedelta(hours=hours, minutes=minutes, seconds=seconds)
    except ValueError as e:
        print(f"Warning: Could not parse timestamp {timestamp_str}: {e}")
        return None

def parse_transcript_line(line):
    """Parse a line from the transcript file"""
    match = re.match(r'\[(.*?) --> (.*?)\] (.*)', line.strip())
    if match:
        start_time = parse_timestamp(match.group(1))
        end_time = parse_timestamp(match.group(2))
        text = match.group(3).strip()
        return start_time, end_time, text
    return None, None, None

def get_slide_timestamps(pdf_path):
    """Extract timestamps from PDF metadata"""
    pdf = fitz.open(pdf_path)
    timestamps = []
    
    # Get timestamps from document metadata
    metadata = pdf.metadata
    if metadata and "subject" in metadata and metadata["subject"]:
        try:
            timestamps = json.loads(metadata["subject"])
            print(f"Found {len(timestamps)} timestamps in PDF metadata")
        except json.JSONDecodeError:
            print("Error: Invalid timestamp metadata in PDF")
    
    pdf.close()
    return timestamps

def get_transcript_for_slide(start_time, end_time, transcript_lines):
    """Get transcript lines that fall between start_time and end_time"""
    matching_lines = []
    
    for line in transcript_lines:
        line_start, line_end, text = parse_transcript_line(line)
        if line_start and line_end and text:
            # Include lines that start during this slide's time window
            if start_time <= line_start and (end_time is None or line_start < end_time):
                matching_lines.append(text)
    
    return '\n'.join(matching_lines) if matching_lines else ""

def create_presentation(pdf_path, transcript_path, output_path):
    # Create a new presentation
    prs = Presentation()
    
    # Read the transcript file
    with open(transcript_path, 'r') as f:
        transcript_lines = [line.strip() for line in f.readlines() if line.strip()]
    
    # Get slide timestamps from PDF
    slide_timestamps = get_slide_timestamps(pdf_path)
    
    if not slide_timestamps:
        print("Error: Could not extract slide timestamps from PDF")
        return
    
    print(f"Found {len(slide_timestamps)} slide timestamps")
    
    # Open the PDF
    pdf = fitz.open(pdf_path)
    
    # Process each page in the PDF
    for page_num in range(len(pdf)):
        print(f"Processing slide {page_num + 1}/{len(pdf)}")
        
        # Get the page
        page = pdf[page_num]
        
        # Convert PDF page to image
        pix = page.get_pixmap()
        img_data = pix.tobytes("png")
        img = Image.open(io.BytesIO(img_data))
        
        # Save the image temporarily
        temp_img_path = f"temp_slide_{page_num}.png"
        img.save(temp_img_path)
        
        # Add a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 is a blank layout
        
        # Add the image to the slide
        left = top = Inches(0)
        pic = slide.shapes.add_picture(temp_img_path, left, top, width=prs.slide_width)
        
        # Get the time window for this slide
        slide_start = parse_timestamp(slide_timestamps[page_num])
        slide_end = parse_timestamp(slide_timestamps[page_num + 1]) if page_num + 1 < len(slide_timestamps) else None
        
        # Add speaker notes from the transcript
        notes_text = get_transcript_for_slide(slide_start, slide_end, transcript_lines)
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
            print(f"Added {len(notes_text.split(chr(10)))} transcript lines to slide {page_num + 1}")
        else:
            print(f"No matching transcript lines for slide {page_num + 1}")
        
        # Remove temporary image
        os.remove(temp_img_path)
    
    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

def main():
    parser = argparse.ArgumentParser(description='Create PowerPoint presentation from slides and transcripts.')
    parser.add_argument('--slides-pdf', default=os.path.join("output", "slides.pdf"), help='Path to slides PDF')
    parser.add_argument('--transcript', default=os.path.join("output", "transcript.txt"), help='Path to transcript file')
    parser.add_argument('--output', default=os.path.join("output", "presentation.pptx"), help='Path to output PowerPoint file')
    args = parser.parse_args()
    
    create_presentation(args.slides_pdf, args.transcript, args.output)

if __name__ == "__main__":
    main() 